import os
import platform
import subprocess
from datetime import datetime
from PIL import Image
from tkinter import simpledialog
from toast import toast_success, toast_warning
from file_manager import MASTER_FOLDER

IMAGE_EXTENSIONS = (".png", ".jpg", ".jpeg")


def get_images():
    """Returns a sorted list of image paths from the capture folder."""
    if not os.path.exists(MASTER_FOLDER):
        return []
    return sorted([
        os.path.join(MASTER_FOLDER, f)
        for f in os.listdir(MASTER_FOLDER)
        if f.lower().endswith(IMAGE_EXTENSIONS)
    ])


def get_comment_for(image_path):
    """Returns the saved comment for an image, or an empty string if none exists."""
    txt_path = os.path.splitext(image_path)[0] + ".txt"
    if os.path.exists(txt_path):
        with open(txt_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    return ""


def get_save_path(settings, extension, parent):
    """
    Builds the output file path from the user's export folder setting.
    If ask_filename is enabled, prompts the user for a custom name.
    Falls back to the script directory if no export folder is configured.
    """
    folder = settings.get("export_folder", "")
    if not folder or not os.path.isdir(folder):
        folder = os.path.dirname(__file__)

    timestamp    = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    default_name = f"Docilio_{timestamp}"

    if settings.get("ask_filename", False):
        user_input = simpledialog.askstring(
            "Save As",
            "Enter a name for your export file:",
            initialvalue=default_name,
            parent=parent
        )
        filename = user_input.strip() if user_input and user_input.strip() else default_name
    else:
        filename = default_name

    return os.path.join(folder, f"{filename}{extension}")


def check_images(parent):
    """Returns the image list, or shows a warning toast and returns None if empty."""
    images = get_images()
    if not images:
        toast_warning("Nothing to export - capture some screenshots first.")
        return None
    return images


def _open_file(path):
    """Opens the exported file in the default application for the current OS."""
    try:
        if platform.system() == "Windows":
            os.startfile(path)
        elif platform.system() == "Darwin":
            subprocess.Popen(["open", path])
        else:
            subprocess.Popen(["xdg-open", path])
    except Exception:
        pass


# WORD EXPORT

def _split_wide_image(img_path):
    """
    Checks whether an image is a wide comparison (aspect ratio >= 2.0) and,
    if so, splits it vertically down the middle into left and right halves.

    Word pages are roughly 6 inches wide. Inserting a side-by-side comparison
    at full page width leaves each half only 3 inches wide, making it hard to
    read. Splitting and stacking the halves gives each one the full page width.

    Returns (left_half, right_half) as PIL Image objects, or None if not wide.
    """
    with Image.open(img_path) as pil_img:
        orig_w, orig_h = pil_img.size

    if (orig_w / orig_h) < 2.0:
        return None

    img   = Image.open(img_path).convert("RGB")
    mid   = img.width // 2
    left  = img.crop((0,   0, mid,       img.height))
    right = img.crop((mid, 0, img.width, img.height))
    return left, right


def _insert_pil_image_to_doc(doc, pil_img, width_inches):
    """
    Inserts a PIL Image into a python-docx Document using an in-memory buffer.
    doc.add_picture() requires a file path or file-like object, so BytesIO is
    used to avoid writing temporary files to disk.
    """
    import io
    from docx.shared import Inches
    buf = io.BytesIO()
    pil_img.save(buf, format="PNG")
    buf.seek(0)
    doc.add_picture(buf, width=Inches(width_inches))


def export_word(settings, parent):
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor

    images = check_images(parent)
    if not images:
        return

    doc = Document()

    for img_path in images:
        halves = _split_wide_image(img_path)

        if halves:
            # Wide comparison image: stack the two halves vertically at full page width
            left_half, right_half = halves
            _insert_pil_image_to_doc(doc, left_half,  width_inches=6)
            doc.add_paragraph("")
            _insert_pil_image_to_doc(doc, right_half, width_inches=6)
        else:
            doc.add_picture(img_path, width=Inches(6))

        comment = get_comment_for(img_path)
        if comment:
            para               = doc.add_paragraph(comment)
            para.style         = doc.styles["Normal"]
            run                = para.runs[0]
            run.font.name      = "Calibri"
            run.font.size      = Pt(10)
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            run.font.italic    = False

        doc.add_paragraph("")

    save_path = get_save_path(settings, ".docx", parent)
    doc.save(save_path)
    toast_success(f"Word document saved - {os.path.basename(save_path)}")
    _open_file(save_path)


# EXCEL EXPORT
#
# openpyxl places images as floating objects anchored to a cell's top-left corner.
# Excel does not auto-resize the row to fit, so images overlap the rows below them.
#
# To fix this, the anchor row height is set in points to match the image height,
# and the row cursor advances by the equivalent number of default-height rows (15pt).
# Conversion: 1px = 0.75pt (since 1pt = 1.33px at 96dpi).

def export_excel(settings, parent):
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage
    from openpyxl.styles import Font, Alignment

    images = check_images(parent)
    if not images:
        return

    wb    = Workbook()
    sheet = wb.active
    sheet.title = "Screenshots"

    DISPLAY_WIDTH_NORMAL = 640   # px, standard screenshots
    DISPLAY_WIDTH_WIDE   = 1200  # px, wide comparison images
    SPACER_ROWS          = 3     # blank rows between each image block
    PT_PER_PX            = 0.75  # 1px = 0.75pt at 96dpi

    current_row = 1
    max_disp_w  = 0

    for img_path in images:
        with Image.open(img_path) as pil_img:
            orig_w, orig_h = pil_img.size

        # Threshold of 2.0 targets panoramic comparison images only.
        # Standard 16:9 screenshots (ratio ~1.78) use the normal width.
        display_w = DISPLAY_WIDTH_WIDE if (orig_w / orig_h) >= 2.0 else DISPLAY_WIDTH_NORMAL
        scale     = display_w / orig_w
        disp_w    = display_w
        disp_h    = int(orig_h * scale)

        max_disp_w = max(max_disp_w, disp_w)

        # Set the image row to the exact pixel height of the image.
        # Excel uses points for row height; 1px = 0.75pt at 96dpi.
        # A single tall row holds the image with no guesswork about row counts.
        row_height_pt = disp_h * PT_PER_PX
        sheet.row_dimensions[current_row].height = row_height_pt

        xl_img        = XLImage(img_path)
        xl_img.width  = disp_w
        xl_img.height = disp_h
        xl_img.anchor = f"A{current_row}"
        sheet.add_image(xl_img)

        # Move to the very next row for the comment - no gap
        current_row += 1

        comment = get_comment_for(img_path)
        if comment:
            cell           = sheet.cell(row=current_row, column=1, value=comment)
            cell.font      = Font(name="Calibri", size=11, color="000000", italic=False)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            sheet.row_dimensions[current_row].height = 30
            current_row += 1

        current_row += SPACER_ROWS

    # Set column A wide enough for the widest image.
    # 1 Excel character unit is roughly 7px at the default font size.
    sheet.column_dimensions["A"].width = (max_disp_w // 7) + 5

    save_path = get_save_path(settings, ".xlsx", parent)
    wb.save(save_path)
    toast_success(f"Excel file saved - {os.path.basename(save_path)}")
    _open_file(save_path)


# PDF EXPORT

def export_pdf(settings, parent):
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.utils import ImageReader
    from reportlab.lib import colors

    images = check_images(parent)
    if not images:
        return

    save_path      = get_save_path(settings, ".pdf", parent)
    page_w, page_h = A4
    margin         = 40

    c = rl_canvas.Canvas(save_path, pagesize=A4)

    for img_path in images:
        with Image.open(img_path) as pil_img:
            orig_w, orig_h = pil_img.size

        # Landscape page for wide images (comparisons, horizontal stitches)
        if orig_w > orig_h:
            c.setPageSize((page_h, page_w))
            p_w, p_h = page_h, page_w
        else:
            c.setPageSize((page_w, page_h))
            p_w, p_h = page_w, page_h

        usable_w = p_w - (margin * 2)
        usable_h = p_h - (margin * 2) - 30

        scale  = usable_w / orig_w
        draw_w = usable_w
        draw_h = orig_h * scale

        if draw_h > usable_h:
            scale  = usable_h / draw_h
            draw_w = draw_w * scale
            draw_h = usable_h

        x_pos = margin + (usable_w - draw_w) / 2
        y_pos = p_h - margin - draw_h

        c.drawImage(
            ImageReader(img_path),
            x_pos, y_pos,
            width=draw_w, height=draw_h,
            preserveAspectRatio=True
        )

        comment = get_comment_for(img_path)
        if comment:
            c.setFont("Helvetica", 9)
            c.setFillColor(colors.black)
            truncated = comment[:160] + ("..." if len(comment) > 160 else "")
            c.drawString(margin, y_pos - 14, truncated)

        c.showPage()

    c.save()
    toast_success(f"PDF saved - {os.path.basename(save_path)}")
    _open_file(save_path)


# POWERPOINT EXPORT

def export_pptx(settings, parent):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    images = check_images(parent)
    if not images:
        return

    prs          = Presentation()
    slide_w      = prs.slide_width
    slide_h      = prs.slide_height
    blank_layout = prs.slide_layouts[6]

    for img_path in images:
        slide = prs.slides.add_slide(blank_layout)

        with Image.open(img_path) as pil_img:
            orig_w, orig_h = pil_img.size

        comment   = get_comment_for(img_path)
        comment_h = Inches(0.4) if comment else 0

        margin = Inches(0.2)
        max_w  = slide_w - (margin * 2)
        max_h  = slide_h - (margin * 2) - comment_h
        scale  = min(max_w / orig_w, max_h / orig_h)

        img_w = int(orig_w * scale)
        img_h = int(orig_h * scale)
        left  = (slide_w - img_w) // 2
        top   = margin

        slide.shapes.add_picture(img_path, left, top, width=img_w, height=img_h)

        if comment:
            txBox        = slide.shapes.add_textbox(
                margin, top + img_h + Inches(0.1),
                slide_w - (margin * 2), comment_h
            )
            tf           = txBox.text_frame
            tf.word_wrap = True
            para         = tf.paragraphs[0]
            run          = para.add_run()
            run.text           = comment
            run.font.name      = "Calibri"
            run.font.size      = Pt(10)
            run.font.italic    = False
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    save_path = get_save_path(settings, ".pptx", parent)
    prs.save(save_path)
    toast_success(f"PowerPoint saved - {os.path.basename(save_path)}")
    _open_file(save_path)